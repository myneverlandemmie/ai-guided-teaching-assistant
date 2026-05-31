"""课次教学材料文本提取工具。

只做教师常见教学材料的基础文本提取，不做 OCR、PDF 解析或版式还原。
"""

from __future__ import annotations

import re
from datetime import date, datetime, time
from pathlib import Path

try:
    from docx import Document
    from docx.table import Table, _Cell
    from docx.text.paragraph import Paragraph
except ImportError:  # pragma: no cover - 缺依赖时由调用方显示错误
    Document = None
    Table = None
    _Cell = None
    Paragraph = None

try:
    from pptx import Presentation
except ImportError:  # pragma: no cover - 缺依赖时由调用方显示错误
    Presentation = None

try:
    from openpyxl import load_workbook
except ImportError:  # pragma: no cover - 缺依赖时由调用方显示错误
    load_workbook = None

SUPPORTED_MATERIAL_SUFFIXES = {".txt", ".md", ".docx", ".pptx", ".xlsx"}
MAX_XLSX_EXTRACTED_CHARS = 120_000
FOOTER_PATTERNS = (
    re.compile(r"©\s*Microsoft Corporation", re.IGNORECASE),
    re.compile(r"All rights reserved", re.IGNORECASE),
    re.compile(r"^\d{4}[/-]\d{1,2}[/-]\d{1,2}(?:\s+\d{1,2}:\d{2}(?::\d{2})?)?$"),
    re.compile(r"^\d{1,2}[/-]\d{1,2}[/-]\d{2,4}(?:\s+\d{1,2}:\d{2}(?::\d{2})?)?$"),
)


class LessonMaterialExtractionError(ValueError):
    """教学材料文本提取失败。"""


def clean_extracted_text(text: str) -> str:
    """清洗提取文本，减少模板噪声但保留教学内容。

    Args:
        text: 从文档中提取的原始文本。

    Returns:
        清理空白、明显页脚、重复空行和整行重复后的文本。

    Raises:
        不主动抛出业务异常。
    """

    cleaned_lines: list[str] = []
    previous_line = ""
    previous_blank = False
    for raw_line in text.splitlines():
        line = re.sub(r"[ \t\f\v]+", " ", raw_line).strip()
        # Word 表格标签常被拆成“重 点”“难 点”，这里仅规范化明确的教学栏目名。
        line = re.sub(r"^重\s*点\s*([:：])?", "重点：", line)
        line = re.sub(r"^难\s*点\s*([:：])?", "难点：", line)
        if any(pattern.search(line) for pattern in FOOTER_PATTERNS):
            continue
        if not line:
            if cleaned_lines and not previous_blank:
                cleaned_lines.append("")
            previous_blank = True
            previous_line = ""
            continue
        if line == previous_line:
            continue
        cleaned_lines.append(line)
        previous_line = line
        previous_blank = False

    return "\n".join(cleaned_lines).strip()


def _iter_docx_blocks(document: object):
    """按 Word 文档主体顺序产出段落和表格。"""

    if Paragraph is None or Table is None:
        return
    for child in document.element.body.iterchildren():
        if child.tag.endswith("}p"):
            yield Paragraph(child, document)
        elif child.tag.endswith("}tbl"):
            yield Table(child, document)


def _extract_docx_cell_text(cell: object) -> str:
    """提取 Word 表格单元格内的多段文本。"""

    paragraphs = [paragraph.text.strip() for paragraph in cell.paragraphs if paragraph.text.strip()]
    return "\n".join(paragraphs).strip()


def _extract_docx_table_text(table: object) -> list[str]:
    """提取 Word 表格文本，按行组织并仅在行内跳过合并单元格重复。"""

    parts: list[str] = []
    for row in table.rows:
        cells: list[str] = []
        seen_cell_xml_in_row: set[str] = set()
        for cell in row.cells:
            cell_key = str(cell._tc.xml)
            if cell_key in seen_cell_xml_in_row:
                continue
            seen_cell_xml_in_row.add(cell_key)
            text = _extract_docx_cell_text(cell)
            if text:
                cells.append(text)
        if not cells:
            continue
        if len(cells) == 2:
            parts.append(f"{cells[0]}：\n{cells[1]}")
        else:
            parts.append("\n---\n".join(cells))
    return parts


def extract_text_from_docx(file_path: str | Path) -> str:
    """从 `.docx` 中提取段落和表格单元格文本。

    Args:
        file_path: `.docx` 文件路径。

    Returns:
        清洗后的文本。

    Raises:
        LessonMaterialExtractionError: 缺少依赖、文件无法解析或未提取到文本。
    """

    if Document is None:
        raise LessonMaterialExtractionError("缺少 python-docx 依赖，无法解析 .docx。")

    try:
        document = Document(str(file_path))
    except Exception as exc:  # noqa: BLE001 - 需要转成教师可理解的错误
        raise LessonMaterialExtractionError(".docx 文本提取失败，请复制 Word 中的文字粘贴到文本框。") from exc

    parts: list[str] = []
    for block in _iter_docx_blocks(document):
        if Paragraph is not None and isinstance(block, Paragraph):
            if block.text.strip():
                parts.append(block.text)
        elif Table is not None and isinstance(block, Table):
            # Word 教案常把教学目标、过程等内容放在表格中，按表格出现位置提取。
            parts.extend(_extract_docx_table_text(block))

    text = clean_extracted_text("\n".join(parts))
    if not text:
        raise LessonMaterialExtractionError("未从 .docx 中提取到可用文本，请复制 Word 中的文字粘贴到文本框。")
    return text


def _extract_text_from_pptx_shape(shape: object) -> list[str]:
    """提取单个 PPT shape 中的文本。"""

    parts: list[str] = []
    if getattr(shape, "has_text_frame", False):
        text = getattr(shape, "text", "").strip()
        if text:
            parts.append(text)
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append("\t".join(cells))
    return parts


def extract_text_from_pptx(file_path: str | Path) -> str:
    """从 `.pptx` 中按幻灯片顺序提取文本框和表格文本。

    Args:
        file_path: `.pptx` 文件路径。

    Returns:
        清洗后的文本，每页包含 `【Slide n】` 标记。

    Raises:
        LessonMaterialExtractionError: 缺少依赖、文件无法解析或未提取到文本。
    """

    if Presentation is None:
        raise LessonMaterialExtractionError("缺少 python-pptx 依赖，无法解析 .pptx。")

    try:
        presentation = Presentation(str(file_path))
    except Exception as exc:  # noqa: BLE001 - 需要转成教师可理解的错误
        raise LessonMaterialExtractionError(".pptx 文本提取失败，请复制 PPT 中的文字粘贴到文本框。") from exc

    slide_blocks: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        slide_parts: list[str] = []
        for shape in slide.shapes:
            slide_parts.extend(_extract_text_from_pptx_shape(shape))
        slide_text = clean_extracted_text("\n".join(slide_parts))
        if slide_text:
            slide_blocks.append(f"【Slide {index}】\n{slide_text}")

    text = clean_extracted_text("\n\n".join(slide_blocks))
    if not text:
        raise LessonMaterialExtractionError("未从 .pptx 中提取到可用文本，请复制 PPT 中的文字粘贴到文本框。")
    return text


def _format_xlsx_cell_value(value: object) -> str:
    """将 openpyxl 单元格值转换为适合进入教学材料的文本。"""

    if value is None:
        return ""
    if isinstance(value, datetime):
        return value.strftime("%Y-%m-%d %H:%M:%S")
    if isinstance(value, date):
        return value.strftime("%Y-%m-%d")
    if isinstance(value, time):
        return value.strftime("%H:%M:%S")
    return str(value).strip()


def _limit_xlsx_text(text: str) -> str:
    """限制大表格提取文本长度，避免单个材料无限增长。"""

    if len(text) <= MAX_XLSX_EXTRACTED_CHARS:
        return text
    clipped = text[:MAX_XLSX_EXTRACTED_CHARS].rstrip()
    return (
        f"{clipped}\n\n"
        "提示：该 XLSX 表格资料内容较长，已按长度限制截取；请教师结合原始文件确认。"
    )


def extract_text_from_xlsx(file_path: str | Path) -> str:
    """从 `.xlsx` 工作簿中提取工作表名和非空行文本。

    Args:
        file_path: `.xlsx` 文件路径。

    Returns:
        Markdown 风格的工作表文本。

    Raises:
        LessonMaterialExtractionError: 缺少依赖、文件无法解析或未提取到文本。
    """

    if load_workbook is None:
        raise LessonMaterialExtractionError("缺少 openpyxl 依赖，无法解析 .xlsx。")

    try:
        workbook = load_workbook(str(file_path), read_only=True, data_only=True)
    except Exception as exc:  # noqa: BLE001 - 需要转成教师可理解的错误
        raise LessonMaterialExtractionError(".xlsx 表格文本提取失败，请检查文件是否损坏或复制表格内容粘贴到文本框。") from exc

    sheet_blocks: list[str] = []
    try:
        for worksheet in workbook.worksheets:
            rows: list[str] = []
            for row in worksheet.iter_rows(values_only=True):
                cells = [_format_xlsx_cell_value(value) for value in row]
                non_empty_cells = [cell for cell in cells if cell]
                if non_empty_cells:
                    rows.append("\t".join(non_empty_cells))
            if rows:
                sheet_blocks.append(f"## Sheet: {worksheet.title}\n\n" + "\n".join(rows))
    finally:
        workbook.close()

    if not sheet_blocks:
        raise LessonMaterialExtractionError("未从 .xlsx 中提取到可用文本，请复制表格内容粘贴到文本框。")

    text = "# XLSX 表格资料提取\n\n" + "\n\n".join(sheet_blocks)
    return _limit_xlsx_text(clean_extracted_text(text))


def extract_text_from_lesson_material(file_path: str | Path, filename: str) -> str:
    """按文件扩展名提取课次材料文本。

    Args:
        file_path: 已保存的上传文件路径。
        filename: 原始文件名，用于判断扩展名。

    Returns:
        清洗后的文本。

    Raises:
        LessonMaterialExtractionError: 文件类型不支持或解析失败。
    """

    path = Path(file_path)
    suffix = Path(filename).suffix.lower()
    if suffix == ".xls":
        raise LessonMaterialExtractionError("暂不支持旧版 .xls 表格文件。请另存为 .xlsx 后上传，或复制表格内容粘贴到文本框。")
    if suffix not in SUPPORTED_MATERIAL_SUFFIXES:
        raise LessonMaterialExtractionError(
            "当前支持粘贴文本或上传 .txt / .md / .docx / .pptx / .xlsx 文件；暂不支持 .xls、PDF、图片、扫描件和旧版 .doc / .ppt。"
        )
    if suffix in {".txt", ".md"}:
        return clean_extracted_text(path.read_text(encoding="utf-8", errors="replace"))
    if suffix == ".docx":
        return extract_text_from_docx(path)
    if suffix == ".pptx":
        return extract_text_from_pptx(path)
    if suffix == ".xlsx":
        return extract_text_from_xlsx(path)
    raise LessonMaterialExtractionError("不支持的教学材料文件类型。")

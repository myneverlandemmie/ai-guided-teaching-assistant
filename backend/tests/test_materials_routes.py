from pathlib import Path

import pytest
from docx import Document
from pptx import Presentation
from sqlalchemy import select

from app import main
from app.models.course_plan import PlannedLesson
from app.models.lesson import LessonMaterial
from app.routes import materials as materials_routes
from tests.support.course_plan_helpers import (
    _build_test_client,
    _create_course,
    _upload_sample_plan,
    anyio_backend,
    inline_threadpool_for_tests,
)


def _create_minimal_docx(path: Path) -> None:
    document = Document()
    document.add_paragraph("教学目标：理解 INNER JOIN")
    table = document.add_table(rows=1, cols=2)
    table.cell(0, 0).text = "重点"
    table.cell(0, 1).text = "GROUP BY 与 HAVING"
    document.save(path)


def _create_table_lesson_plan_docx(path: Path) -> None:
    document = Document()
    document.add_paragraph("条件查询教学设计")
    table = document.add_table(rows=7, cols=2)
    rows = [
        ("教学目标", "知识目标\n掌握WHERE子句的使用方法\n能力目标\n能够根据任务选择合适的筛选条件\n素养目标\n形成规范书写 SQL 的习惯"),
        ("重点", "WHERE、AND、OR、NOT 条件筛选"),
        ("难点", "多个条件组合时的逻辑关系"),
        ("教学过程", "组织教学\n复习导入\n新课讲授：WHERE 条件查询\n课堂练习：查询指定类别商品"),
        ("课堂练习", "使用 IN关键字 完成多条件筛选。"),
        ("布置作业", "完成 WHERE 子句练习题。"),
        ("教学反思", "学生对多条件组合还需要更多示例。"),
    ]
    for index, (label, value) in enumerate(rows):
        table.cell(index, 0).text = label
        table.cell(index, 1).text = value
    document.save(path)


def _create_minimal_pptx(path: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    title = slide.shapes.title
    title.text = "SQL 高级查询"
    text_box = slide.shapes.add_textbox(0, 0, 5000000, 1000000)
    text_box.text_frame.text = "INNER JOIN 与 OUTER JOIN"
    footer = slide.shapes.add_textbox(0, 1000000, 5000000, 1000000)
    footer.text_frame.text = "© Microsoft Corporation All rights reserved"
    presentation.save(path)


@pytest.mark.anyio
async def test_can_submit_text_lesson_material(tmp_path: Path) -> None:
    client, session_factory = _build_test_client(tmp_path)
    course = _create_course(session_factory)
    try:
        await _upload_sample_plan(client, course)
        with session_factory() as session:
            selected_id = session.scalar(select(PlannedLesson.id).order_by(PlannedLesson.id))
            assert selected_id is not None
        await client.post(
            "/course-plan-uploads/1/confirm",
            data={"planned_lesson_ids": str(selected_id)},
            follow_redirects=False,
        )

        response = await client.post(
            "/lessons/1/materials",
            data={
                "title": "本节课导入材料",
                "material_type": "pasted_text",
                "content": "SELECT 与 WHERE 的课堂说明。",
            },
            follow_redirects=False,
        )

        assert response.status_code == 303
        assert response.headers["location"] == "/lessons/1"
        with session_factory() as session:
            materials = session.scalars(select(LessonMaterial)).all()
            assert len(materials) == 1
            assert materials[0].lesson_id == 1
            assert materials[0].title == "本节课导入材料"
            assert materials[0].content == "SELECT 与 WHERE 的课堂说明。"
            assert materials[0].file_path is None
    finally:
        await client.aclose()
        main.app.dependency_overrides.clear()


@pytest.mark.anyio
async def test_lesson_detail_shows_saved_material(tmp_path: Path) -> None:
    client, session_factory = _build_test_client(tmp_path)
    course = _create_course(session_factory)
    try:
        await _upload_sample_plan(client, course)
        with session_factory() as session:
            selected_id = session.scalar(select(PlannedLesson.id).order_by(PlannedLesson.id))
            assert selected_id is not None
        await client.post(
            "/course-plan-uploads/1/confirm",
            data={"planned_lesson_ids": str(selected_id)},
            follow_redirects=False,
        )
        await client.post(
            "/lessons/1/materials",
            data={
                "title": "课堂讲义文本",
                "material_type": "pasted_text",
                "content": "本节课讲解简单查询和条件筛选。",
            },
            follow_redirects=False,
        )

        response = await client.get("/lessons/1")

        assert response.status_code == 200
        assert "课堂讲义文本" in response.text
        assert "本节课讲解简单查询和条件筛选。" in response.text
    finally:
        await client.aclose()
        main.app.dependency_overrides.clear()


@pytest.mark.anyio
async def test_lesson_detail_shows_material_support_scope(tmp_path: Path) -> None:
    client, session_factory = _build_test_client(tmp_path)
    course = _create_course(session_factory)
    try:
        await _upload_sample_plan(client, course)
        with session_factory() as session:
            selected_id = session.scalar(select(PlannedLesson.id).order_by(PlannedLesson.id))
            assert selected_id is not None
        await client.post(
            "/course-plan-uploads/1/confirm",
            data={"planned_lesson_ids": str(selected_id)},
            follow_redirects=False,
        )

        response = await client.get("/lessons/1")

        assert response.status_code == 200
        assert "添加到本课次" in response.text
        assert "txt / md / docx / pptx / xlsx" in response.text
        assert "xlsx 将提取表格文本" in response.text
        assert "暂不支持 xls、PDF、图片或扫描件" in response.text
    finally:
        await client.aclose()
        main.app.dependency_overrides.clear()


@pytest.mark.anyio
async def test_unsupported_lesson_material_file_type_shows_hint(tmp_path: Path) -> None:
    client, session_factory = _build_test_client(tmp_path)
    course = _create_course(session_factory)
    try:
        await _upload_sample_plan(client, course)
        with session_factory() as session:
            selected_id = session.scalar(select(PlannedLesson.id).order_by(PlannedLesson.id))
            assert selected_id is not None
        await client.post(
            "/course-plan-uploads/1/confirm",
            data={"planned_lesson_ids": str(selected_id)},
            follow_redirects=False,
        )

        response = await client.post(
            "/lessons/1/materials",
            data={"title": "PDF 材料", "material_type": "supplementary", "content": ""},
            files={"files": ("lesson.pdf", b"fake pdf", "application/pdf")},
        )

        assert response.status_code == 400
        assert response.status_code != 500
        assert "暂不支持该文件类型" in response.text
        assert "txt / md / docx / pptx / xlsx" in response.text
        assert "暂不支持 xls、PDF、图片或扫描件" in response.text
        with session_factory() as session:
            assert session.scalars(select(LessonMaterial)).all() == []
    finally:
        await client.aclose()
        main.app.dependency_overrides.clear()


@pytest.mark.anyio
async def test_oversized_lesson_material_file_shows_hint_without_saving(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    monkeypatch.setattr(materials_routes, "MAX_LESSON_MATERIAL_UPLOAD_BYTES", 8)
    client, session_factory = _build_test_client(tmp_path)
    course = _create_course(session_factory)
    try:
        await _upload_sample_plan(client, course)
        with session_factory() as session:
            selected_id = session.scalar(select(PlannedLesson.id).order_by(PlannedLesson.id))
            assert selected_id is not None
        await client.post(
            "/course-plan-uploads/1/confirm",
            data={"planned_lesson_ids": str(selected_id)},
            follow_redirects=False,
        )

        response = await client.post(
            "/lessons/1/materials",
            data={"title": "超限资料", "material_type": "supplementary", "content": ""},
            files={"files": ("too-large.txt", b"012345678", "text/plain")},
        )

        assert response.status_code == 400
        assert response.status_code != 500
        assert "文件过大，请拆分资料后上传。" in response.text
        with session_factory() as session:
            assert session.scalars(select(LessonMaterial)).all() == []
        upload_dir = tmp_path / "lesson-materials"
        assert not upload_dir.exists() or list(upload_dir.iterdir()) == []
    finally:
        await client.aclose()
        main.app.dependency_overrides.clear()


@pytest.mark.anyio
async def test_v2_material_upload_error_stays_on_materials_outline_page(tmp_path: Path) -> None:
    client, session_factory = _build_test_client(tmp_path)
    course = _create_course(session_factory)
    try:
        await _upload_sample_plan(client, course)
        with session_factory() as session:
            selected_id = session.scalar(select(PlannedLesson.id).order_by(PlannedLesson.id))
            assert selected_id is not None
        await client.post(
            "/course-plan-uploads/1/confirm",
            data={"planned_lesson_ids": str(selected_id)},
            follow_redirects=False,
        )

        response = await client.post(
            "/lessons/1/materials",
            data={
                "title": "PDF 材料",
                "input_mode": "file_upload",
                "material_category": "supplementary",
                "return_to": "/ui-v2/lessons/1/materials-outline",
            },
            files={"files": ("lesson.pdf", b"fake pdf", "application/pdf")},
        )

        assert response.status_code == 400
        assert response.status_code != 500
        assert "lesson-materials-v2" in response.text
        assert "课程资料整理" in response.text
        assert "数据库应用与数据分析" in response.text
        assert "上传课次资料" in response.text
        assert "暂不支持该文件类型" in response.text
        assert "课次详情" not in response.text
        with session_factory() as session:
            assert session.scalars(select(LessonMaterial)).all() == []
    finally:
        await client.aclose()
        main.app.dependency_overrides.clear()


@pytest.mark.anyio
async def test_can_upload_docx_lesson_material_with_paragraph_and_table_text(tmp_path: Path) -> None:
    client, session_factory = _build_test_client(tmp_path)
    course = _create_course(session_factory)
    docx_path = tmp_path / "lesson.docx"
    _create_minimal_docx(docx_path)
    try:
        await _upload_sample_plan(client, course)
        with session_factory() as session:
            selected_id = session.scalar(select(PlannedLesson.id).order_by(PlannedLesson.id))
            assert selected_id is not None
        await client.post(
            "/course-plan-uploads/1/confirm",
            data={"planned_lesson_ids": str(selected_id)},
            follow_redirects=False,
        )

        response = await client.post(
            "/lessons/1/materials",
            data={"title": "Word 教案", "material_type": "lesson_plan", "content": ""},
            files={
                "files": (
                    "lesson.docx",
                    docx_path.read_bytes(),
                    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                )
            },
            follow_redirects=False,
        )

        assert response.status_code == 303
        with session_factory() as session:
            material = session.scalar(select(LessonMaterial))
            assert material is not None
            assert "教学目标：理解 INNER JOIN" in material.content
            assert "重点" in material.content
            assert "GROUP BY 与 HAVING" in material.content
            assert material.file_path is not None

        page_response = await client.get("/lessons/1")
        assert page_response.status_code == 200
        assert "lesson.docx" in page_response.text
        assert str(tmp_path) not in page_response.text
        assert "data/uploads" not in page_response.text
    finally:
        await client.aclose()
        main.app.dependency_overrides.clear()


@pytest.mark.anyio
async def test_can_extract_table_style_docx_lesson_plan_completely(tmp_path: Path) -> None:
    client, session_factory = _build_test_client(tmp_path)
    course = _create_course(session_factory)
    docx_path = tmp_path / "table-lesson-plan.docx"
    _create_table_lesson_plan_docx(docx_path)
    try:
        await _upload_sample_plan(client, course)
        with session_factory() as session:
            selected_id = session.scalar(select(PlannedLesson.id).where(PlannedLesson.lesson_code == "0402"))
            assert selected_id is not None
        await client.post(
            "/course-plan-uploads/1/confirm",
            data={"planned_lesson_ids": str(selected_id)},
            follow_redirects=False,
        )

        response = await client.post(
            "/lessons/1/materials",
            data={"title": "", "material_type": "lesson_plan", "content": ""},
            files={
                "files": (
                    "table-lesson-plan.docx",
                    docx_path.read_bytes(),
                    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                )
            },
            follow_redirects=False,
        )

        assert response.status_code == 303
        with session_factory() as session:
            material = session.scalar(select(LessonMaterial))
            assert material is not None
            assert material.title == "0402-教案"
            assert "知识目标" in material.content
            assert "掌握WHERE子句的使用方法" in material.content
            assert "能力目标" in material.content
            assert "素养目标" in material.content
            assert "重点" in material.content
            assert "难点" in material.content
            assert "教学过程" in material.content
            assert "课堂练习" in material.content
            assert "教学反思" in material.content
            assert "\n" in material.content
            assert len(material.content) > 180
            assert material.content.count("教学过程") == 1

        page_response = await client.get("/lessons/1")
        assert page_response.status_code == 200
        assert "table-lesson-plan.docx" in page_response.text
        assert str(tmp_path) not in page_response.text
        assert "data/uploads" not in page_response.text
        assert "掌握WHERE子句的使用方法" in page_response.text
    finally:
        await client.aclose()
        main.app.dependency_overrides.clear()


@pytest.mark.anyio
async def test_default_lesson_material_title_adds_sequence_for_duplicates(tmp_path: Path) -> None:
    client, session_factory = _build_test_client(tmp_path)
    course = _create_course(session_factory)
    docx_path = tmp_path / "table-lesson-plan.docx"
    _create_table_lesson_plan_docx(docx_path)
    try:
        await _upload_sample_plan(client, course)
        with session_factory() as session:
            selected_id = session.scalar(select(PlannedLesson.id).where(PlannedLesson.lesson_code == "0402"))
            assert selected_id is not None
        await client.post(
            "/course-plan-uploads/1/confirm",
            data={"planned_lesson_ids": str(selected_id)},
            follow_redirects=False,
        )

        for _ in range(2):
            response = await client.post(
                "/lessons/1/materials",
                data={"title": "", "material_type": "lesson_plan", "content": ""},
                files={
                    "files": (
                        "table-lesson-plan.docx",
                        docx_path.read_bytes(),
                        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                    )
                },
                follow_redirects=False,
            )
            assert response.status_code == 303

        with session_factory() as session:
            titles = session.scalars(select(LessonMaterial.title).order_by(LessonMaterial.id)).all()
            assert titles == ["0402-教案", "0402-教案（2）"]
    finally:
        await client.aclose()
        main.app.dependency_overrides.clear()


@pytest.mark.anyio
async def test_can_upload_pptx_lesson_material_with_slide_text_and_cleaned_footer(tmp_path: Path) -> None:
    client, session_factory = _build_test_client(tmp_path)
    course = _create_course(session_factory)
    pptx_path = tmp_path / "lesson.pptx"
    _create_minimal_pptx(pptx_path)
    try:
        await _upload_sample_plan(client, course)
        with session_factory() as session:
            selected_id = session.scalar(select(PlannedLesson.id).order_by(PlannedLesson.id))
            assert selected_id is not None
        await client.post(
            "/course-plan-uploads/1/confirm",
            data={"planned_lesson_ids": str(selected_id)},
            follow_redirects=False,
        )

        response = await client.post(
            "/lessons/1/materials",
            data={"title": "PPT 材料", "material_type": "course_ppt", "content": ""},
            files={
                "files": (
                    "lesson.pptx",
                    pptx_path.read_bytes(),
                    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                )
            },
            follow_redirects=False,
        )

        assert response.status_code == 303
        with session_factory() as session:
            material = session.scalar(select(LessonMaterial))
            assert material is not None
            assert "【Slide 1】" in material.content
            assert "SQL 高级查询" in material.content
            assert "INNER JOIN 与 OUTER JOIN" in material.content
            assert "© Microsoft Corporation" not in material.content
            assert "All rights reserved" not in material.content
            assert material.file_path is not None
    finally:
        await client.aclose()
        main.app.dependency_overrides.clear()


@pytest.mark.anyio
async def test_can_upload_multiple_lesson_material_files(tmp_path: Path) -> None:
    client, session_factory = _build_test_client(tmp_path)
    course = _create_course(session_factory)
    guide_path = tmp_path / "guide.md"
    note_path = tmp_path / "note.txt"
    guide_path.write_text("实训指导：完成 INNER JOIN 查询", encoding="utf-8")
    note_path.write_text("补充资料：GROUP BY 练习", encoding="utf-8")
    try:
        await _upload_sample_plan(client, course)
        with session_factory() as session:
            selected_id = session.scalar(select(PlannedLesson.id).order_by(PlannedLesson.id))
            assert selected_id is not None
        await client.post(
            "/course-plan-uploads/1/confirm",
            data={"planned_lesson_ids": str(selected_id)},
            follow_redirects=False,
        )

        response = await client.post(
            "/lessons/1/materials",
            data={"title": "批量资料", "material_type": "supplementary", "content": ""},
            files=[
                ("files", ("guide.md", guide_path.read_bytes(), "text/markdown")),
                ("files", ("note.txt", note_path.read_bytes(), "text/plain")),
            ],
            follow_redirects=False,
        )

        assert response.status_code == 303
        with session_factory() as session:
            materials = session.scalars(select(LessonMaterial).order_by(LessonMaterial.id)).all()
            assert len(materials) == 2
            assert "实训指导：完成 INNER JOIN 查询" in materials[0].content
            assert "补充资料：GROUP BY 练习" in materials[1].content
    finally:
        await client.aclose()
        main.app.dependency_overrides.clear()


@pytest.mark.anyio
async def test_docx_extraction_deduplicates_repeated_lines(tmp_path: Path) -> None:
    client, session_factory = _build_test_client(tmp_path)
    course = _create_course(session_factory)
    docx_path = tmp_path / "duplicate.docx"
    document = Document()
    document.add_paragraph("教学目标：理解连接查询")
    document.add_paragraph("教学目标：理解连接查询")
    table = document.add_table(rows=2, cols=1)
    table.cell(0, 0).text = "教学过程：演示 INNER JOIN"
    table.cell(1, 0).text = "教学过程：演示 INNER JOIN"
    document.save(docx_path)
    try:
        await _upload_sample_plan(client, course)
        with session_factory() as session:
            selected_id = session.scalar(select(PlannedLesson.id).order_by(PlannedLesson.id))
            assert selected_id is not None
        await client.post(
            "/course-plan-uploads/1/confirm",
            data={"planned_lesson_ids": str(selected_id)},
            follow_redirects=False,
        )

        response = await client.post(
            "/lessons/1/materials",
            data={"title": "去重教案", "material_type": "lesson_plan", "content": ""},
            files={
                "files": (
                    "duplicate.docx",
                    docx_path.read_bytes(),
                    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                )
            },
            follow_redirects=False,
        )

        assert response.status_code == 303
        with session_factory() as session:
            material = session.scalar(select(LessonMaterial))
            assert material is not None
            assert material.content.count("教学目标：理解连接查询") == 1
            assert material.content.count("教学过程：演示 INNER JOIN") == 1
    finally:
        await client.aclose()
        main.app.dependency_overrides.clear()


@pytest.mark.anyio
async def test_can_delete_lesson_material(tmp_path: Path) -> None:
    client, session_factory = _build_test_client(tmp_path)
    course = _create_course(session_factory)
    try:
        await _upload_sample_plan(client, course)
        with session_factory() as session:
            selected_id = session.scalar(select(PlannedLesson.id).order_by(PlannedLesson.id))
            assert selected_id is not None
        await client.post(
            "/course-plan-uploads/1/confirm",
            data={"planned_lesson_ids": str(selected_id)},
            follow_redirects=False,
        )
        await client.post(
            "/lessons/1/materials",
            data={"title": "待删除资料", "material_type": "pasted_text", "content": "传错的资料"},
            follow_redirects=False,
        )
        with session_factory() as session:
            material = session.scalar(select(LessonMaterial))
            assert material is not None
            material_id = material.id

        response = await client.post(f"/lesson-materials/{material_id}/delete", follow_redirects=False)

        assert response.status_code == 303
        assert response.headers["location"] == "/lessons/1"
        with session_factory() as session:
            assert session.scalars(select(LessonMaterial)).all() == []
    finally:
        await client.aclose()
        main.app.dependency_overrides.clear()
